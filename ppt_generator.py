from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
from datetime import datetime
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
import subprocess
import platform
import shutil

class PPTGenerator:
    """PowerPoint 및 보고서 생성 클래스"""
    
    def __init__(self):
        self.prs = None
    
    def create_presentation(self, company, project_name, year, month, photos_path):
        """PPT 프레젠테이션 생성"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # 슬라이드 1: 제목
        self._add_title_slide(company, project_name, year, month)
        
        # 슬라이드 2: 프로젝트 개요
        self._add_overview_slide(company, project_name, year, month)
        
        # 슬라이드 3-5: 사진 슬라이드 (작업전/중/후)
        self._add_photo_slides(photos_path)
        
        # 슬라이드 6: 요약
        self._add_summary_slide()
        
        return self.prs
    
    def _add_title_slide(self, company, project_name, year, month):
        """제목 슬라이드 추가"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 빈 슬라이드
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = company
        title_frame.paragraphs[0].font.size = Pt(60)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True
        
        # 부제목
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"{project_name}\n{year}년 {month}월"
        subtitle_frame.paragraphs[0].font.size = Pt(32)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.word_wrap = True
    
    def _add_overview_slide(self, company, project_name, year, month):
        """프로젝트 개요 슬라이드"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "📋 프로젝트 정보"
        
        # 내용
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        info = [
            f"회사: {company}",
            f"현장명: {project_name}",
            f"년도: {year}년",
            f"월: {month}월",
            f"작성일: {datetime.now().strftime('%Y년 %m월 %d일')}"
        ]
        
        for i, line in enumerate(info):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(24)
            p.level = 0
            p.space_before = Pt(6)
    
    def _add_photo_slides(self, photos_path):
        """사진 슬라이드 추가 (작업전/중/후)"""
        photos_path = Path(photos_path)
        categories = ['작업전', '작업중', '작업후']
        
        for category in categories:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = f"📷 작업 사진 - {category}"
            
            category_path = photos_path / category
            
            if category_path.exists():
                photos = sorted([f for f in category_path.glob('*') if f.is_file()])
                
                if not photos:
                    # 사진이 없는 경우
                    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
                    text_frame = content_box.text_frame
                    text_frame.text = f"📁 {category} 사진이 없습니다."
                    text_frame.paragraphs[0].font.size = Pt(28)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                else:
                    # 최대 3개 사진 표시
                    for idx, photo in enumerate(photos[:3]):
                        try:
                            left = Inches(0.5 + idx * 3.1)
                            top = Inches(2)
                            slide.shapes.add_picture(str(photo), left, top, width=Inches(2.8))
                        except Exception as e:
                            print(f"사진 추가 오류 ({photo.name}): {e}")
            else:
                # 폴더 자체가 없는 경우
                content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
                text_frame = content_box.text_frame
                text_frame.text = f"📁 {category} 폴더가 없습니다."
                text_frame.paragraphs[0].font.size = Pt(28)
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_summary_slide(self):
        """요약 슬라이드"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "✅ 작업 완료"
        
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(3))
        text_frame = content_box.text_frame
        text_frame.text = "📊 보고서 작성이 완료되었습니다."
        text_frame.paragraphs[0].font.size = Pt(36)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.word_wrap = True
    
    def save_ppt(self, file_path):
        """PPT 파일 저장"""
        if self.prs is None:
            raise ValueError("프레젠테이션이 생성되지 않았습니다.")
        self.prs.save(str(file_path))
        return file_path
    
    def ppt_to_pdf(self, ppt_path, pdf_path):
        """PPT를 PDF로 변환 (LibreOffice 사용)"""
        try:
            pdf_path = Path(pdf_path)
            ppt_path = Path(ppt_path)
            
            # LibreOffice 명령어
            if platform.system() == 'Windows':
                # Windows
                soffice_paths = [
                    r'C:\Program Files\LibreOffice\program\soffice.exe',
                    r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
                ]
                soffice_path = None
                for path in soffice_paths:
                    if Path(path).exists():
                        soffice_path = path
                        break
                
                if not soffice_path:
                    print("⚠️ LibreOffice를 찾을 수 없습니다. PDF 변환 스킵합니다.")
                    return None
                
                subprocess.run([
                    soffice_path,
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', str(pdf_path.parent),
                    str(ppt_path)
                ], capture_output=True, timeout=60)
            
            elif platform.system() == 'Darwin':
                # macOS
                subprocess.run([
                    '/usr/local/bin/libreoffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', str(pdf_path.parent),
                    str(ppt_path)
                ], capture_output=True, timeout=60)
            
            else:
                # Linux
                subprocess.run([
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', str(pdf_path.parent),
                    str(ppt_path)
                ], capture_output=True, timeout=60)
            
            # PDF 파일이 생성되었는지 확인
            generated_pdf = pdf_path.parent / f"{ppt_path.stem}.pdf"
            if generated_pdf.exists():
                # 원하는 이름으로 복사
                if generated_pdf != pdf_path:
                    shutil.move(str(generated_pdf), str(pdf_path))
                return pdf_path
            else:
                print(f"⚠️ PDF 변환 실패: {generated_pdf}")
                return None
        
        except subprocess.TimeoutExpired:
            print("⚠️ PDF 변환 시간 초과")
            return None
        except Exception as e:
            print(f"⚠️ PDF 변환 오류: {e}")
            return None
    
    def generate_excel(self, company, project_name, year, month, photos_info, file_path):
        """엑셀 파일 생성"""
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "보고서"
            
            # 헤더 스타일
            header_fill = PatternFill(start_color="2C3E50", end_color="2C3E50", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=12)
            
            # 제목
            ws['A1'] = f"{company} - {project_name} ({year}년 {month}월)"
            ws['A1'].font = Font(bold=True, size=14, color="2C3E50")
            ws.merge_cells('A1:D1')
            ws.row_dimensions[1].height = 25
            
            # 프로젝트 정보
            ws['A3'] = "📋 프로젝트 정보"
            ws['A3'].font = header_font
            ws['A3'].fill = header_fill
            ws.row_dimensions[3].height = 20
            
            row = 4
            info = [
                ("회사", company),
                ("현장명", project_name),
                ("년도", year),
                ("월", month),
                ("작성일", datetime.now().strftime('%Y-%m-%d'))
            ]
            
            for label, value in info:
                ws[f'A{row}'] = label
                ws[f'A{row}'].font = Font(bold=True)
                ws[f'B{row}'] = value
                row += 1
            
            # 사진 통계
            row += 2
            ws[f'A{row}'] = "📷 사진 통계"
            ws[f'A{row}'].font = header_font
            ws[f'A{row}'].fill = header_fill
            ws.row_dimensions[row].height = 20
            
            row += 1
            ws[f'A{row}'] = "카테고리"
            ws[f'B{row}'] = "개수"
            
            for cell in [ws[f'A{row}'], ws[f'B{row}']]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="3498DB", end_color="3498DB", fill_type="solid")
                cell.alignment = Alignment(horizontal="center", vertical="center")
            
            row += 1
            total_photos = 0
            for category, count in photos_info.items():
                ws[f'A{row}'] = category
                ws[f'B{row}'] = count
                ws[f'B{row}'].alignment = Alignment(horizontal="center")
                total_photos += count
                row += 1
            
            # 합계
            ws[f'A{row}'] = "합계"
            ws[f'A{row}'].font = Font(bold=True)
            ws[f'B{row}'] = total_photos
            ws[f'B{row}'].alignment = Alignment(horizontal="center")
            ws[f'B{row}'].font = Font(bold=True)
            
            # 열 너비 조정
            ws.column_dimensions['A'].width = 20
            ws.column_dimensions['B'].width = 15
            
            wb.save(str(file_path))
            return file_path
        except Exception as e:
            print(f"엑셀 생성 오류: {e}")
            raise
